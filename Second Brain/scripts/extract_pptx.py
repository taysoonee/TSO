from pptx import Presentation
import sys

def extract_text(path):
    prs = Presentation(path)
    for i, slide in enumerate(prs.slides):
        print(f"--- Slide {i+1} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                print(shape.text)

if __name__ == "__main__":
    extract_text(sys.argv[1])
